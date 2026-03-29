"""
Render PPTX slides to PNG images.

Priority:
  1. PowerPoint COM (Windows, if available) — pixel-perfect
  2. python-pptx + Pillow — reads real slide content, draws title/bullets/colors
"""

import io
import os
import base64
import uuid
import logging
import tempfile
import textwrap
from typing import Optional

try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PILLOW = True
except ImportError:
    HAS_PILLOW = False

logger = logging.getLogger(__name__)

# ── Colour Palette (matches the Neon theme) ────────────────────────────────────
BG_COLOR   = (10, 15, 30)       # #0A0F1E  deep navy
PANEL_BG   = (18, 28, 52)       # header / bullet area
ACCENT     = (245, 83, 61)      # #F5533D  neon red
TEXT_LIGHT = (241, 245, 249)    # #F1F5F9
TEXT_DIM   = (148, 163, 184)    # #94A3B8
WHITE      = (255, 255, 255)


def _find_font(size: int) -> ImageFont.FreeTypeFont:
    """Try common system font paths; fall back to Pillow default."""
    candidates = [
        "arial.ttf",
        "Arial.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/calibri.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _rgb_from_pptx(color_obj) -> tuple:
    """Safely extract RGB tuple from a python-pptx color object."""
    try:
        from pptx.dml.color import RGBColor
        rgb = color_obj.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return TEXT_LIGHT


def _extract_slide_content(slide) -> tuple[str, list[str], Optional[bytes]]:
    """Extract (title, [bullet1, bullet2, ...], image_bytes) from a python-pptx slide object."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    title = ""
    bullets: list[str] = []
    image_bytes = None

    for shape in slide.shapes:
        # Check for image
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            if not image_bytes:
                try:
                    image_bytes = shape.image.blob
                except Exception:
                    pass

        if not getattr(shape, "has_text_frame", False):
            continue

        text = shape.text_frame.text.strip()
        if not text:
            continue

        # Shape name heuristics + placeholder index
        try:
            ph = shape.placeholder_format
            if ph.idx == 0:
                title = text
                continue
        except Exception:
            pass

        # Everything else is a content/bullet block
        for para in shape.text_frame.paragraphs:
            line = para.text.strip()
            # Ignore bullet/spacer placeholders like "." or "•"
            if line and line != title and len(line) > 1 and line not in ["•", "-", ".", "*", "o", "°"]:
                bullets.append(line)

    # If no explicit title placeholder found, use first shape text
    if not title:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip()
                break

    return title, bullets[:8], image_bytes   # cap at 8 bullets


def _render_slide_pillow(slide, slide_index: int, width: int, height: int) -> str:
    """Render one python-pptx slide to a base64 PNG using Pillow."""
    # ── Fonts ──────────────────────────────────────────────────────────────────
    font_num      = _find_font(max(14, height // 28))    # slide number
    font_title    = _find_font(max(20, height // 20))    # title
    font_bullet   = _find_font(max(14, height // 40))    # bullet text
    font_tiny     = _find_font(max(10, height // 70))    # label

    img  = Image.new("RGB", (width, height), BG_COLOR)
    draw = ImageDraw.Draw(img)

    pad = int(width * 0.04)       # 4% horizontal padding
    top_bar_h = int(height * 0.14)

    # ── Top header band ────────────────────────────────────────────────────────
    draw.rectangle([0, 0, width, top_bar_h], fill=PANEL_BG)

    # Accent left stripe
    stripe_w = max(4, int(width * 0.004))
    draw.rectangle([0, 0, stripe_w, top_bar_h], fill=ACCENT)

    # Slide index badge
    badge_size = int(top_bar_h * 0.55)
    badge_x    = pad
    badge_y    = (top_bar_h - badge_size) // 2
    draw.rectangle([badge_x, badge_y, badge_x + badge_size, badge_y + badge_size], fill=ACCENT)
    num_str = f"{slide_index + 1:02d}"
    draw.text(
        (badge_x + badge_size // 2, badge_y + badge_size // 2),
        num_str, fill=WHITE, anchor="mm", font=font_num,
    )

    # Title text
    title, bullets, image_bytes = _extract_slide_content(slide)
    title_x = badge_x + badge_size + int(pad * 0.6)
    title_y = (top_bar_h - top_bar_h // 3) // 2
    # Truncate title if too long
    max_title_chars = int((width - title_x - pad) / max(1, font_title.size * 0.6))
    display_title = title[:max_title_chars] + ("…" if len(title) > max_title_chars else "")
    draw.text((title_x, title_y), display_title, fill=TEXT_LIGHT, font=font_title)

    # Thin accent line under header
    draw.rectangle([0, top_bar_h, width, top_bar_h + 2], fill=ACCENT)

    # ── Bullet content area ────────────────────────────────────────────────────
    bullet_char = "▸"
    content_y   = top_bar_h + int(height * 0.06)
    line_h      = int(height * 0.09)
    max_bullet_width = width - (pad * 2) - int(width * 0.04)

    # Estimate ~chars per line
    approx_char_w = max(6, font_bullet.size // 2)
    # If image is present, the text area is much narrower
    actual_bullet_width = (max_bullet_width // 2) if image_bytes else max_bullet_width
    chars_per_line = max(20, actual_bullet_width // approx_char_w)

    for i, bullet in enumerate(bullets):
        y = content_y + i * line_h
        if y + line_h > height - int(height * 0.06):
            break

        # Bullet dot / arrow
        draw.text((pad, y + 4), bullet_char, fill=ACCENT, font=font_bullet)

        # Bullet text (wrapped)
        wrapped = textwrap.wrap(bullet, width=chars_per_line)
        first_line = wrapped[0] if wrapped else bullet
        draw.text((pad + int(width * 0.035), y), first_line, fill=TEXT_DIM, font=font_bullet)

    # ── Draw Image if present ──────────────────────────────────────────────────
    if image_bytes:
        try:
            img_io = io.BytesIO(image_bytes)
            overlay = Image.open(img_io).convert("RGBA")
            
            # Target bounds for image on right side:
            w = int(width * 0.40)
            h = int(height * 0.72)
            x = int(width * 0.55)
            y = int(height * 0.14)
            
            # Maintain aspect ratio to fit bounds
            overlay.thumbnail((w, h), Image.Resampling.LANCZOS)
            cx = x + (w - overlay.width) // 2
            cy = y + (h - overlay.height) // 2
            
            img.paste(overlay, (cx, cy), overlay)
        except Exception as e:
            logger.warning("Failed to overlay image in Pillow fallback: %s", e)

    # ── Bottom strip ───────────────────────────────────────────────────────────
    bottom_y = height - int(height * 0.045)
    draw.rectangle([0, bottom_y, width, height], fill=PANEL_BG)
    draw.rectangle([0, bottom_y, width, bottom_y + 1], fill=ACCENT)
    draw.text(
        (width - pad, bottom_y + int(height * 0.02)),
        "NEO PPT", fill=ACCENT, anchor="rm", font=font_tiny,
    )

    # ── Encode ─────────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    return f"data:image/png;base64,{b64}"


def render_slides_to_images(
    pptx_buffer: io.BytesIO, width: int = 1920, height: int = 1080
) -> list[Optional[str]]:
    """
    Render every slide in *pptx_buffer* as a base64-encoded PNG.

    Strategy
    --------
    1. Try PowerPoint COM automation (Windows, pixel-perfect).
    2. Fall back to python-pptx + Pillow (cross-platform, content-accurate).
    """
    from pptx import Presentation

    # Parse the presentation once to get slide objects for fallback
    pptx_buffer.seek(0)
    try:
        prs = Presentation(pptx_buffer)
    except Exception as e:
        logger.error("Cannot open PPTX buffer: %s", e)
        return []
    pptx_buffer.seek(0)

    # ── Attempt 1: PowerPoint COM ──────────────────────────────────────────────
    tmp_dir   = tempfile.mkdtemp(prefix="neo_ppt_render_")
    tmp_pptx  = os.path.join(tmp_dir, f"render_{uuid.uuid4().hex[:8]}.pptx")
    slide_dir = os.path.join(tmp_dir, "slides")
    os.makedirs(slide_dir, exist_ok=True)

    try:
        import comtypes.client
        import pythoncom

        with open(tmp_pptx, "wb") as f:
            f.write(pptx_buffer.read())
        pptx_buffer.seek(0)


        pythoncom.CoInitialize()
        images: list[Optional[str]] = []
        try:
            ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
            opened  = ppt_app.Presentations.Open(
                os.path.abspath(tmp_pptx),
                ReadOnly=True,
                WithWindow=False,
            )
            for i, slide in enumerate(opened.Slides):
                out_path = os.path.join(slide_dir, f"slide_{i}.png")
                slide.Export(out_path, "PNG", width, height)
                if os.path.exists(out_path):
                    with open(out_path, "rb") as img_file:
                        b64 = base64.b64encode(img_file.read()).decode("utf-8")
                        images.append(f"data:image/png;base64,{b64}")
                    os.remove(out_path)
                else:
                    # COM slide failed individually — use Pillow for this one
                    images.append(_render_slide_pillow(prs.slides[i], i, width, height) if HAS_PILLOW else None)

            opened.Close()
            ppt_app.Quit()
            logger.info("COM rendered %d slides.", len(images))
            return images
        finally:
            pythoncom.CoUninitialize()
            _cleanup(tmp_dir, tmp_pptx, slide_dir)

    except Exception as e:
        logger.warning("COM rendering unavailable (%s); using Pillow fallback.", e)
        pptx_buffer.seek(0)

    # ── Attempt 2: Pillow content renderer ────────────────────────────────────
    if not HAS_PILLOW:
        logger.error("Pillow not installed — cannot render previews.")
        return [None] * len(prs.slides)

    images = []
    for i, slide in enumerate(prs.slides):
        try:
            images.append(_render_slide_pillow(slide, i, width, height))
        except Exception as e:
            logger.warning("Pillow render failed for slide %d: %s", i, e)
            images.append(None)

    logger.info("Pillow rendered %d slides.", len(images))
    return images


def _cleanup(tmp_dir: str, tmp_pptx: str, slide_dir: str):
    """Best-effort cleanup of temp files."""
    try:
        if os.path.exists(tmp_pptx):
            os.remove(tmp_pptx)
        if os.path.exists(slide_dir):
            for f in os.listdir(slide_dir):
                os.remove(os.path.join(slide_dir, f))
            os.rmdir(slide_dir)
        if os.path.exists(tmp_dir):
            os.rmdir(tmp_dir)
    except Exception:
        pass
