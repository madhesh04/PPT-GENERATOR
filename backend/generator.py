import os
import copy
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Template path ──────────────────────────────────────────────────────────────
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "template.pptx")

from shared.themes import THEME_DATA

# ── Themes ─────────────────────────────────────────────────────────────────────
class Theme:
    def __init__(self, name, main, accent, text, slate, white, surface):
        self.name = name
        self.main = RGBColor(*main)
        self.accent = RGBColor(*accent)
        self.text = RGBColor(*text)
        self.slate = RGBColor(*slate)
        self.white = RGBColor(*white)
        self.surface = RGBColor(*surface)

THEMES = {
    k: Theme(
        v["name"], v["main"], v["accent"], v["text"], 
        v["slate"], v["white"], v["surface"]
    ) for k, v in THEME_DATA.items()
}

DEFAULT_THEME = THEMES["neon"]

SLIDE_W = Inches(20.00)
SLIDE_H = Inches(11.25)


# ── Slide duplication (within same prs — preserves image relationships) ────────

def _duplicate_slide_within(prs: Presentation, source_idx: int):
    """
    Clone slide[source_idx] and append it to the end of prs.
    Works within the SAME Presentation so all embedded image/media parts
    and their relationship IDs are preserved.
    """
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # 1. Copy relationships, building old-rId → new-rId map
    rId_map: dict[str, str] = {}
    for rel in source.part.rels.values():
        if "slideLayout" in rel.reltype or "slideMaster" in rel.reltype:
            continue
        if rel.is_external:
            continue
        new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rId_map[rel.rId] = new_rId

    # 2. Deep-copy the shape tree and remap all rId attribute values
    src_copy = copy.deepcopy(source.shapes._spTree)
    for elem in src_copy.iter():
        for attr, val in list(elem.attrib.items()):
            if val in rId_map:
                elem.attrib[attr] = rId_map[val]

    # 3. Replace the new slide's shape tree
    dst = new_slide.shapes._spTree
    for child in list(dst):
        dst.remove(child)
    for child in src_copy:
        dst.append(child)

    return new_slide


# ── Drawing helpers ────────────────────────────────────────────────────────────

def _add_text_box(slide, text, left, top, width, height,
                  font_size=24, bold=False,
                  color: RGBColor = None,
                  align=PP_ALIGN.LEFT,
                  italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide decorators ───────────────────────────────────────────────────────────

def _decorate_closing_slide(slide, title: str, theme: Theme):
    """Renders a 'Thank You' closing slide."""
    _add_text_box(slide, "Thank You",
                  Inches(1.5), Inches(3.6), Inches(12), Inches(2.5),
                  font_size=72, bold=True, color=theme.white)
    _add_text_box(slide, title,
                  Inches(1.5), Inches(5.8), Inches(14), Inches(1.0),
                  font_size=24, color=theme.main, italic=True)


def _decorate_agenda_slide(slide, slide_titles: list, theme: Theme):
    """Renders a clean 'Agenda' overview slide on the template layout."""
    _add_text_box(slide, "Agenda",
                  Inches(1.0), Inches(0.4), Inches(15.0), Inches(1.0),
                  font_size=36, bold=True, color=theme.text)

    num_items = len(slide_titles)
    
    # Decide layout: single column if 7 or fewer, otherwise two columns
    if num_items <= 7:
        # Single column
        for row, item in enumerate(slide_titles):
            y = 2.0 + (row * 0.9) # Slightly more spacing for fewer items
            _add_text_box(slide, f"{row+1:02d}. {item}",
                          Inches(1.5), Inches(y), Inches(15.0), Inches(0.65),
                          font_size=26, color=theme.slate)
    else:
        # Two columns (split as evenly as possible)
        midpoint = (num_items + 1) // 2
        col1 = slide_titles[:midpoint]
        col2 = slide_titles[midpoint:]

        # Column 1 (Left)
        for row, item in enumerate(col1):
            y = 2.0 + (row * 0.9)
            _add_text_box(slide, f"{row+1:02d}. {item}",
                          Inches(1.5), Inches(y), Inches(8.5), Inches(0.65),
                          font_size=26, color=theme.slate)

        # Column 2 (Right)
        for row, item in enumerate(col2):
            y = 2.0 + (row * 0.9)
            _add_text_box(slide, f"{row+midpoint+1:02d}. {item}",
                          Inches(10.5), Inches(y), Inches(8.5), Inches(0.65),
                          font_size=26, color=theme.slate)



def _decorate_title_slide(slide, title: str, theme: Theme):
    # Only map the title text over the existing template.
    _add_text_box(slide, title,
                  Inches(11.0), Inches(5.8), Inches(8.0), Inches(3.0),
                  font_size=60, bold=True, color=theme.text)


def _decorate_content_slide(slide, slide_title: str, content: list,
                             slide_num: int, theme: Theme, image_bytes: bytes | None = None):
    """Renders a standard content slide with 5 bullets and optional image. No code."""
    has_image = image_bytes is not None

    # Slide Title (full width)
    _add_text_box(slide, slide_title,
                  Inches(1.0), Inches(0.4), Inches(18.0), Inches(1.0),
                  font_size=36, bold=True, color=theme.text)

    # Bullet width: narrower if image present
    bullet_text_width = Inches(9.0) if has_image else Inches(17.0)
    bullet_font_size  = 22 if has_image else 24

    # Bullet Points (up to 5)
    for i, point in enumerate(content[:5]):
        y = 1.8 + (i * 1.5)
        _add_text_box(slide, "•",
                      Inches(1.0), Inches(y - 0.05), Inches(0.5), Inches(1.0),
                      font_size=28, color=theme.main)
        _add_text_box(slide, point,
                      Inches(1.5), Inches(y), bullet_text_width, Inches(1.4),
                      font_size=bullet_font_size, color=theme.slate)

    # Image — right half
    if has_image:
        try:
            img_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(
                img_stream,
                left=Inches(11.0), top=Inches(1.6),
                width=Inches(8.0), height=Inches(8.2)
            )
        except Exception as e:
            import logging
            logging.getLogger(__name__).warning(
                "Failed to add image to slide %d: %s", slide_num, e
            )


def _decorate_code_slide(slide, parent_title: str, code: str,
                          language: str | None, theme: Theme):
    """Renders a full dedicated code slide with light background and dark text."""
    # Title
    display_title = f"Code — {parent_title}"
    _add_text_box(slide, display_title,
                  Inches(1.0), Inches(0.4), Inches(16.0), Inches(0.9),
                  font_size=30, bold=True, color=theme.text)

    # Language badge (top-right)
    if language:
        _add_text_box(slide, language.upper(),
                      Inches(17.0), Inches(0.5), Inches(2.0), Inches(0.4),
                      font_size=12, bold=True,
                      color=theme.main, align=PP_ALIGN.RIGHT)

    # Light code background box (nearly full slide)
    code_bg = slide.shapes.add_shape(
        1, Inches(1.0), Inches(1.5), Inches(18.0), Inches(8.5)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    code_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    code_bg.line.width = Pt(1)

    # Code text — black Consolas on light background
    code_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.8), Inches(17.0), Inches(7.8)
    )
    tf = code_box.text_frame
    tf.word_wrap = True

    code_lines = code.replace("\\n", "\n").split("\n")
    for line_idx, line in enumerate(code_lines):
        if line_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)
        run.font.name = "Consolas"
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        run.font.bold = False


# ── Main entry point ───────────────────────────────────────────────────────────

def create_presentation(title: str, slide_data: list,
                         image_bytes_list: list | None = None,
                         theme_name: str = "neon"
                        ) -> tuple[io.BytesIO, str]:
    """
    Generate a PPTX in memory. Returns (BytesIO, filename).
    Code slides are rendered as separate full slides right after their parent.
    """
    theme = THEMES.get(theme_name.lower(), DEFAULT_THEME)

    if os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
        n = len(slide_data)

        # Count code slides — each gets a dedicated full slide
        code_slide_count = sum(
            1 for sc in slide_data
            if sc.get("code") and str(sc["code"]).strip()
            and str(sc.get("code", "")).lower() not in ("null", "none")
        )
        total_content_slots = n + code_slide_count

        # We need: Title (1) + Agenda (1) + Content (total_content_slots) + Closing (1)
        # Standard template has 2 slides: [Title, Content]

        # 1. Create Agenda slot
        _duplicate_slide_within(prs, 1)

        # 2. Create Content + Code slots
        for _ in range(max(0, total_content_slots - 1)):
            _duplicate_slide_within(prs, 1)

        # 3. Create Closing slot
        _duplicate_slide_within(prs, 0)
        closing_slide_idx = len(prs.slides) - 1

        # Decorate
        _decorate_title_slide(prs.slides[0], title, theme)

        titles = [s.get("title", f"Slide {i+1}") for i, s in enumerate(slide_data)]
        _decorate_agenda_slide(prs.slides[1], titles, theme)

        # Slides 2+: Content (and code slides interleaved)
        prs_idx = 2
        for i, sc in enumerate(slide_data):
            img = image_bytes_list[i] if image_bytes_list and i < len(image_bytes_list) else None

            # Content slide (bullets + optional image, NO code)
            _decorate_content_slide(
                prs.slides[prs_idx],
                slide_title=sc.get("title", f"Slide {i+1}"),
                content=sc.get("content", []),
                slide_num=i + 1,
                theme=theme,
                image_bytes=img,
            )
            prs_idx += 1

            # Dedicated code slide right after, if present
            code_val = sc.get("code")
            if code_val and str(code_val).strip() and str(code_val).lower() not in ("null", "none"):
                _decorate_code_slide(
                    prs.slides[prs_idx],
                    parent_title=sc.get("title", f"Slide {i+1}"),
                    code=str(code_val),
                    language=sc.get("language"),
                    theme=theme,
                )
                prs_idx += 1

        # Last Slide: Closing
        _decorate_closing_slide(prs.slides[closing_slide_idx], title, theme)
    else:
        raise FileNotFoundError(f"Template file missing: {TEMPLATE_PATH}. A template is required for generation.")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)
    filename = f"{safe.replace(' ', '_')}.pptx"
    return buf, filename

