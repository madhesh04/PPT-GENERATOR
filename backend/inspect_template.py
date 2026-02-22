from pptx import Presentation
from pptx.util import Pt

prs = Presentation("template.pptx")

for s_idx, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE [{s_idx}] — layout: '{slide.slide_layout.name}'")
    print(f"  Total shapes: {len(slide.shapes)}")
    for i, shape in enumerate(slide.shapes):
        print(f"\n  Shape [{i}] name='{shape.name}' shape_type={shape.shape_type}")
        print(f"    left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
        if shape.has_text_frame:
            tf = shape.text_frame
            print(f"    has_text_frame=True  paragraphs={len(tf.paragraphs)}")
            for p_idx, para in enumerate(tf.paragraphs):
                text = para.text
                if text.strip():
                    font_size = None
                    bold = None
                    if para.runs:
                        font_size = para.runs[0].font.size
                        bold = para.runs[0].font.bold
                        if font_size:
                            font_size = round(font_size / 12700)  # convert EMU to pt
                    print(f"      para[{p_idx}] text='{text[:80]}' font_size={font_size}pt bold={bold}")
        else:
            print(f"    has_text_frame=False")
